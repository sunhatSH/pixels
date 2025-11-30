#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将 PIXELS_LAMBDA_PROJECT_SUMMARY_PPT.md 内容填充到 PowerPoint 模板（增强版）
支持插入图表和更好的格式控制
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys

TEMPLATE_PATH = "/Users/sunhao/Desktop/实验室/中国人民大学模板-22.pptx"
OUTPUT_PATH = "/Users/sunhao/Desktop/实验室/Pixels_Lambda_项目总结_完整版.pptx"
CHARTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "pptdata")

def add_text_to_shape(shape, text, bold=False, font_size=18, alignment=PP_ALIGN.LEFT):
    """向形状添加文本"""
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # 处理多行文本
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.alignment = alignment
        run = p.add_run()
        run.text = line.strip() if line.strip() else ""
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 处理项目符号
        if line.strip().startswith('•') or line.strip().startswith('✓') or line.strip().startswith('-'):
            p.level = 0
        elif line.strip().startswith('  -'):
            p.level = 1

def insert_image(slide, image_path, left, top, width, height):
    """在幻灯片中插入图片"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                    width=Inches(width), height=Inches(height))
            return True
        except Exception as e:
            print(f"⚠️ 无法插入图片 {image_path}: {e}")
            return False
    else:
        print(f"⚠️ 图片文件不存在: {image_path}")
        return False

def create_summary_slides(prs):
    """创建项目总结幻灯片"""
    
    # 获取模板的第一个布局（标题页）
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    # 幻灯片 1: 标题页
    slide = prs.slides.add_slide(title_layout)
    title_shape = slide.shapes.title if slide.shapes.title else None
    subtitle_shape = None
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            subtitle_shape = shape
            break
    
    if title_shape:
        add_text_to_shape(title_shape, "Pixels Lambda Worker 项目总结", bold=True, font_size=32)
    if subtitle_shape:
        add_text_to_shape(subtitle_shape, "基于 AWS Lambda 的 Serverless 数据处理系统\n\n中国人民大学", bold=False, font_size=18)
    
    # 幻灯片 2: 目录
    slide = prs.slides.add_slide(content_layout)
    title_shape = slide.shapes.title if slide.shapes.title else None
    
    if title_shape:
        add_text_to_shape(title_shape, "目录", bold=True, font_size=28)
    
    # 创建内容文本框
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    items = [
        "1. Lambda 和 Invoker 工作协作流程",
        "2. 从编码到测试、再到获取性能数据的流程",
        "3. 测试文件信息（大小、结构）",
        "4. 测试结果"
    ]
    
    for i, item in enumerate(items):
        p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(24)
        run.font.bold = True
        if i < len(items) - 1:
            p.space_after = Pt(12)
    
    # 幻灯片 3: 架构概览
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "Pixels-Turbo 架构概览", bold=True, font_size=28)
    
    content = """核心组件：

• Coordinator (本地/EC2)
  - Planner: 生成物理执行计划
  - Trino: SQL 查询引擎  
  - Invoker: 调用 Lambda 的客户端

• AWS Lambda (云端)
  - Worker: 执行实际数据处理
  - 按需启动、自动扩展

• AWS S3 (对象存储)
  - 输入数据文件 (.pxl)
  - 输出结果文件 (.pxl)

数据流：Coordinator → Invoker → Lambda → S3"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=18)
    
    # 幻灯片 4: 完整请求流程
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "完整请求流程（端到端）", bold=True, font_size=28)
    
    content = """步骤 1: Coordinator 生成计划
  ↓
步骤 2: Invoker 序列化并调用 AWS Lambda
  ↓
步骤 3: Lambda Worker 执行（冷启动 ~100ms，热启动 ~10ms）
  ↓
步骤 4: S3 读取数据 → 内存处理 → S3 写入结果
  ↓
步骤 5: 返回结果并协调下一步

关键特点：
• 异步调用（CompletableFuture<Output>）
• 支持并发多个 Worker
• 自动处理 AWS SDK 网络通信
• 端到端耗时：1-10 秒（取决于数据大小）"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 幻灯片 5: 开发与部署流程
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "开发与部署流程", bold=True, font_size=28)
    
    content = """自动化部署流程（9 个步骤）：

1. 本地编码 (Mac)
2. Git 提交与推送
3. EC2 编译 (Maven)
4. 下载 JAR 到本地
5. 上传 JAR 到 S3
6. 创建/更新 Lambda 函数
7. 调用 Lambda 测试
8. 从 CloudWatch Logs 提取性能数据
9. 生成 CSV 报告

工具：
• auto-deploy.sh: 自动化部署脚本
• test-all-lambda-workers.sh: 测试脚本
• download-csv-metrics.py: 性能数据提取"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 幻灯片 6: 测试文件信息
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "测试文件信息", bold=True, font_size=28)
    
    content = """S3 测试数据文件：

• large_test_data.pxl: 240.2 MiB
  - ScanWorker 主测试文件
  - 列式存储格式 (.pxl)
  
• example.pxl: 790 Bytes
• input.pxl: 790 Bytes

Pixels 文件特点：
✓ 列式存储，压缩高效（2-10x 压缩比）
✓ 支持选择性列读取（列投影）
✓ 支持行组（Row Group）级别过滤
✓ 包含完整的 Schema 元数据"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 尝试插入文件大小对比图
    chart_path = os.path.join(CHARTS_DIR, "chart5_file_sizes.png")
    if os.path.exists(chart_path):
        insert_image(slide, chart_path, 5.8, 2.5, 4, 3)
    
    # 幻灯片 7: Lambda Workers 部署状态
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "Lambda Workers 部署状态", bold=True, font_size=28)
    
    content = """已部署的 Lambda 函数（总计：9 个）：

✓ ScanWorker
✓ PartitionWorker
✓ AggregationWorker
✓ BroadcastJoinWorker
✓ PartitionedJoinWorker
✓ SortedJoinWorker
✓ BroadcastChainJoinWorker
✓ PartitionedChainJoinWorker
✓ SortWorker

部署结果：
• 所有函数均已成功部署
• 所有函数均可正常调用
• 已创建 CloudWatch Log Groups"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 尝试插入部署状态图
    chart_path = os.path.join(CHARTS_DIR, "chart3_workers_deployment.png")
    if os.path.exists(chart_path):
        insert_image(slide, chart_path, 5.8, 2.5, 4, 3)
    
    # 幻灯片 8: 测试结果
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "ScanWorker 测试结果", bold=True, font_size=28)
    
    content = """测试执行结果：

✓ Lambda 调用成功
✓ 数据读取完成
✓ 数据处理完成
✓ 结果写入 S3
✓ 性能指标记录

测试输入：
• 输入文件: large_test_data.pxl (240.2 MiB)
• 列投影: 3 列
• 过滤器: 空（无过滤）
• 输出路径: s3://home-sunhao/output/

测试环境：
• Lambda 内存: 4096 MB
• Lambda 超时: 15 分钟
• 区域: us-east-2"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 尝试插入测试结果图
    chart_path = os.path.join(CHARTS_DIR, "chart4_test_results.png")
    if os.path.exists(chart_path):
        insert_image(slide, chart_path, 5.8, 2.5, 4, 3)
    
    # 幻灯片 9: 性能指标结果
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "性能指标结果", bold=True, font_size=28)
    
    content = """四阶段性能指标（测试数据：240.2 MiB）：

• READ: 9354 ms (26.19%)
  从 S3 读取数据

• COMPUTE: 9718 ms (27.21%)
  过滤、投影、编码

• WRITE_CACHE: 13110 ms (36.71%)
  写入 Lambda 内存缓存

• WRITE_FILE: 3533 ms (9.89%)
  持久化到 S3

总耗时: 35715 ms (约 35.7 秒)
内存使用: 3068 MB / 4096 MB (74.9%)"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 尝试插入性能指标图
    chart_path = os.path.join(CHARTS_DIR, "chart1_performance_timing.png")
    if os.path.exists(chart_path):
        insert_image(slide, chart_path, 5.8, 2, 4.2, 3.5)
    
    # 幻灯片 10: 性能占比分析
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "性能占比分析", bold=True, font_size=28)
    
    content = """性能分析：

存储 I/O 占比：
• S3 Storage (READ + WRITE_FILE): 36.08%
• Total storage (read+write): 72.79%

时间分类：
• 存储 I/O: 12887 ms (36.08%)
• 计算操作: 9718 ms (27.21%)
• 内存操作: 13110 ms (36.71%)

关键发现：
✓ 存储 I/O 是主要瓶颈（36.08%）
✓ WRITE_CACHE 占比最高（36.71%）
✓ 计算时间合理（包含数据编码）"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 尝试插入占比分析图
    chart_path = os.path.join(CHARTS_DIR, "chart2_performance_percentage.png")
    if os.path.exists(chart_path):
        insert_image(slide, chart_path, 5.8, 2, 4.2, 3.5)
    
    # 幻灯片 11: 执行时间线
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "执行时间线", bold=True, font_size=28)
    
    # 尝试插入时间线图
    chart_path = os.path.join(CHARTS_DIR, "chart7_execution_timeline.png")
    if os.path.exists(chart_path):
        insert_image(slide, chart_path, 1, 2.5, 8, 4)
    else:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
        add_text_to_shape(textbox, "执行顺序：READ → COMPUTE → WRITE_CACHE → WRITE_FILE\n总耗时: 35715 ms", 
                         bold=False, font_size=18)
    
    # 幻灯片 12: 完成的工作总结
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "完成的工作总结", bold=True, font_size=28)
    
    content = """✓ 学习并理解了 Lambda 和 Invoker 的协作流程
  Coordinator → Invoker → AWS Lambda → Worker → S3
  完整的异步请求响应机制

✓ 实现了从编码到测试的完整自动化流程
  Git 同步 → EC2 编译 → S3 部署 → Lambda 更新 → 测试执行
  一键部署脚本自动化整个流程

✓ 验证了测试文件的有效性
  240.2 MiB 测试文件成功处理

✓ 获取并分析了性能数据
  四阶段性能指标成功记录，CSV 报告自动生成

✓ 部署了 9 个 Lambda Workers
  所有函数均可正常调用，已创建 CloudWatch Log Groups"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)
    
    # 幻灯片 13: 下一步工作
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "下一步工作", bold=True, font_size=28)
    
    content = """• 为其他 Workers 准备正确的测试输入
  - Partition、Join、Aggregation 等 Workers 需要特定的输入格式
  - 使用 ScanWorker 的输出作为其他 Worker 的输入

• 端到端测试
  - Scan → Partition → Join → Aggregation 完整流程
  - 验证完整的数据处理管道

• 性能优化
  - 根据性能数据分析，优化瓶颈阶段
  - 当前存储 I/O 占比 36.08%，是主要优化目标
  - 考虑使用 S3 Transfer Acceleration 或缓存机制"""
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    add_text_to_shape(textbox, content, bold=False, font_size=16)

def main():
    print("=" * 60)
    print("开始填充 PowerPoint 模板（增强版）...")
    print("=" * 60)
    
    try:
        # 打开模板
        print(f"📖 读取模板: {TEMPLATE_PATH}")
        prs = Presentation(TEMPLATE_PATH)
        print(f"✓ 模板已加载，包含 {len(prs.slides)} 个幻灯片")
        
        # 创建总结幻灯片
        print("\n📝 创建项目总结幻灯片...")
        create_summary_slides(prs)
        
        # 保存
        print(f"\n💾 保存到: {OUTPUT_PATH}")
        prs.save(OUTPUT_PATH)
        
        print("=" * 60)
        print("✅ PowerPoint 文件生成完成！")
        print(f"📁 输出文件: {OUTPUT_PATH}")
        print(f"📊 包含 {len(prs.slides)} 个幻灯片")
        print("=" * 60)
        
    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()

if __name__ == '__main__':
    main()

