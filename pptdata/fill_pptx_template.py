#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将 PIXELS_LAMBDA_PROJECT_SUMMARY_PPT.md 内容填充到 PowerPoint 模板
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys

# 添加父目录到路径以读取 markdown
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

TEMPLATE_PATH = "/Users/sunhao/Desktop/实验室/中国人民大学模板-22.pptx"
OUTPUT_PATH = "/Users/sunhao/Desktop/实验室/Pixels_Lambda_项目总结.pptx"

def add_text_to_shape(shape, text, bold=False, font_size=18):
    """向形状添加文本"""
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)

def create_summary_slides(prs):
    """创建项目总结幻灯片"""
    
    # 幻灯片 1: 标题页
    slide = prs.slides[0] if len(prs.slides) > 0 else prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = None
    subtitle_shape = None
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "标题" in shape.text_frame.text or shape == slide.shapes.title:
                title_shape = shape
            else:
                subtitle_shape = shape
    
    if title_shape:
        add_text_to_shape(title_shape, "Pixels Lambda Worker 项目总结", bold=True, font_size=32)
    if subtitle_shape:
        add_text_to_shape(subtitle_shape, "基于 AWS Lambda 的 Serverless 数据处理系统", bold=False, font_size=18)
    
    # 添加新幻灯片：目录
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title if slide.shapes.title else None
    content_shape = None
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            content_shape = shape
            break
    
    if title_shape:
        add_text_to_shape(title_shape, "目录", bold=True, font_size=24)
    
    if content_shape:
        content_shape.text_frame.clear()
        p = content_shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        items = [
            "1. Lambda 和 Invoker 工作协作流程",
            "2. 从编码到测试、再到获取性能数据的流程",
            "3. 测试文件信息（大小、结构）",
            "4. 测试结果"
        ]
        
        for item in items:
            run = p.add_run()
            run.text = item + "\n"
            run.font.size = Pt(20)
            run.font.bold = False
    
    # 幻灯片 3: Lambda 和 Invoker 工作协作流程 - 架构
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "Pixels-Turbo 架构概览", bold=True, font_size=24)
    
    content = """核心组件：

• Coordinator (本地/EC2)
  - Planner: 生成物理执行计划
  - Trino: SQL 查询引擎
  - Invoker: 调用 Lambda 的客户端

• AWS Lambda (云端)
  - Worker: 执行实际数据处理
  - 按需启动、自动扩展

• AWS S3 (对象存储)
  - 输入数据文件 (.pxl)
  - 输出结果文件 (.pxl)"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 4: 完整请求流程
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "完整请求流程（端到端）", bold=True, font_size=24)
    
    content = """步骤 1: Coordinator 生成计划
  ↓
步骤 2: Invoker 序列化并调用 AWS Lambda
  ↓
步骤 3: Lambda Worker 执行
  ↓
步骤 4: S3 读取数据 → 内存处理 → S3 写入结果
  ↓
步骤 5: 返回结果并协调下一步

特点：
• 异步调用（CompletableFuture）
• 支持并发多个 Worker
• 自动处理 AWS SDK 网络通信"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 5: 开发与部署流程
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "开发与部署流程", bold=True, font_size=24)
    
    content = """1. 本地编码 (Mac)
2. Git 提交与推送
3. EC2 编译 (Maven)
4. 下载 JAR 到本地
5. 上传 JAR 到 S3
6. 创建/更新 Lambda 函数
7. 调用 Lambda 测试
8. 从 CloudWatch Logs 提取性能数据
9. 生成 CSV 报告"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=18)
    
    # 幻灯片 6: 测试文件信息
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "测试文件信息", bold=True, font_size=24)
    
    content = """S3 测试数据文件：

• large_test_data.pxl: 240.2 MiB
  - ScanWorker 主测试文件
  - 列式存储格式
  
• example.pxl: 790 Bytes
• input.pxl: 790 Bytes

文件特点：
✓ 列式存储，压缩高效
✓ 支持选择性列读取
✓ 支持行组级别过滤"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 7: Lambda Workers 部署状态
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "Lambda Workers 部署状态", bold=True, font_size=24)
    
    content = """已部署的 Lambda 函数（总计：9 个）：

✓ ScanWorker
✓ PartitionWorker
✓ AggregationWorker
✓ BroadcastJoinWorker
✓ PartitionedJoinWorker
✓ SortedJoinWorker
✓ BroadcastChainJoinWorker
✓ PartitionedChainJoinWorker
✓ SortWorker

所有函数均已成功部署并可调用！"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 8: 测试结果
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "测试结果", bold=True, font_size=24)
    
    content = """ScanWorker 测试结果：

✓ Lambda 调用成功
✓ 数据读取完成
✓ 数据处理完成
✓ 结果写入 S3
✓ 性能指标记录

测试输入：
• 输入文件: large_test_data.pxl (240.2 MiB)
• 列投影: 3 列
• 过滤器: 空（无过滤）"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 9: 性能指标结果
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "性能指标结果", bold=True, font_size=24)
    
    content = """四阶段性能指标：

• READ: 9354 ms (26.19%)
  从 S3 读取数据

• COMPUTE: 9718 ms (27.21%)
  过滤、投影、编码

• WRITE_CACHE: 13110 ms (36.71%)
  写入 Lambda 内存缓存

• WRITE_FILE: 3533 ms (9.89%)
  持久化到 S3

总耗时: 35715 ms (约 35.7 秒)
内存使用: 3068 MB / 4096 MB"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 10: 其他 Workers 测试状态
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "其他 Workers 测试状态", bold=True, font_size=24)
    
    content = """测试结果总结：

✓ ScanWorker: 成功执行

⚠️ 其他 8 个 Workers: 函数正常，需要正确的输入参数
  - PartitionWorker: 需要 tableInfo
  - AggregationWorker: 需要 aggregationInfo
  - Join Workers: 需要 leftTable/smallTable
  - SortWorker: 需要 tableInfo

结论：
• 所有 Lambda 函数都可以成功调用
• 所有函数都创建了 CloudWatch Log Groups
• 失败原因都是输入参数不完整，不是函数本身的问题"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=14)
    
    # 幻灯片 11: 完成的工作总结
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "完成的工作总结", bold=True, font_size=24)
    
    content = """✓ 学习并理解了 Lambda 和 Invoker 的协作流程
  Coordinator → Invoker → AWS Lambda → Worker → S3

✓ 实现了从编码到测试的完整自动化流程
  Git 同步 → EC2 编译 → S3 部署 → Lambda 更新 → 测试执行

✓ 验证了测试文件的有效性
  240.2 MiB 测试文件成功处理

✓ 获取并分析了性能数据
  四阶段性能指标成功记录，CSV 报告自动生成

✓ 部署了 9 个 Lambda Workers
  所有函数均可正常调用"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)
    
    # 幻灯片 12: 下一步工作
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        add_text_to_shape(slide.shapes.title, "下一步工作", bold=True, font_size=24)
    
    content = """• 为其他 Workers 准备正确的测试输入
  - Partition、Join、Aggregation 等 Workers 需要特定的输入格式

• 端到端测试
  - Scan → Partition → Join → Aggregation 完整流程

• 性能优化
  - 根据性能数据分析，优化瓶颈阶段
  - 当前存储 I/O 占比 36.08%，是主要瓶颈"""
    
    if content_shape:
        add_text_to_shape(content_shape, content, bold=False, font_size=16)

def main():
    print("=" * 60)
    print("开始填充 PowerPoint 模板...")
    print("=" * 60)
    
    try:
        # 打开模板
        print(f"📖 读取模板: {TEMPLATE_PATH}")
        prs = Presentation(TEMPLATE_PATH)
        print(f"✓ 模板已加载，包含 {len(prs.slides)} 个幻灯片")
        
        # 创建总结幻灯片
        print("\n📝 创建项目总结幻灯片...")
        create_summary_slides(prs)
        
        # 保存
        print(f"\n💾 保存到: {OUTPUT_PATH}")
        prs.save(OUTPUT_PATH)
        
        print("=" * 60)
        print("✅ PowerPoint 文件生成完成！")
        print(f"📁 输出文件: {OUTPUT_PATH}")
        print("=" * 60)
        
    except Exception as e:
        print(f"❌ 错误: {e}")
        import traceback
        traceback.print_exc()

if __name__ == '__main__':
    main()

